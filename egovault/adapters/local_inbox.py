"""LocalInboxAdapter — ingests a local file drop folder."""
from __future__ import annotations

import logging
from datetime import datetime, timezone
from pathlib import Path
from typing import TYPE_CHECKING, Iterator

from egovault.core.adapter import BasePlatformAdapter
from egovault.core.registry import register
from egovault.core.schema import NormalizedRecord
from egovault.utils.hashing import compute_content_hash, compute_file_id

if TYPE_CHECKING:
    from egovault.core.store import VaultStore

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Single source of truth for all supported file extensions.
# Each entry: suffix → (record_type, mime_type)
# Adding a new format requires ONE entry here only.
# ---------------------------------------------------------------------------
_FILE_TYPE_INFO: dict[str, tuple[str, str]] = {
    # Documents / notes
    ".md":   ("note",        "text/markdown"),
    ".txt":  ("note",        "text/plain"),
    ".html": ("note",        "text/html"),
    ".pdf":  ("document",    "application/pdf"),
    ".docx": ("document",    "application/vnd.openxmlformats-officedocument.wordprocessingml.document"),
    ".epub": ("document",    "application/epub+zip"),
    ".pptx": ("document",    "application/vnd.openxmlformats-officedocument.presentationml.presentation"),
    ".xlsx": ("spreadsheet", "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"),
    # Source code
    ".ino":   ("code", "text/plain"),
    ".c":     ("code", "text/x-csrc"),
    ".cpp":   ("code", "text/x-c++src"),
    ".h":     ("code", "text/x-chdr"),
    ".hpp":   ("code", "text/x-c++hdr"),
    ".py":    ("code", "text/x-python"),
    ".js":    ("code", "text/javascript"),
    ".ts":    ("code", "text/typescript"),
    ".jsx":   ("code", "text/jsx"),
    ".tsx":   ("code", "text/tsx"),
    ".java":  ("code", "text/x-java-source"),
    ".cs":    ("code", "text/x-csharp"),
    ".rs":    ("code", "text/x-rust"),
    ".go":    ("code", "text/x-go"),
    ".rb":    ("code", "text/x-ruby"),
    ".php":   ("code", "text/x-php"),
    ".swift": ("code", "text/x-swift"),
    ".kt":    ("code", "text/x-kotlin"),
    ".sh":    ("code", "text/x-sh"),
    ".bash":  ("code", "text/x-sh"),
    ".ps1":   ("code", "text/plain"),
    ".bat":   ("code", "text/plain"),
    # Config / data
    ".json":  ("code", "application/json"),
    ".yaml":  ("code", "text/yaml"),
    ".yml":   ("code", "text/yaml"),
    ".toml":  ("code", "text/toml"),
    ".xml":   ("code", "text/xml"),
    ".csv":   ("code", "text/csv"),
    ".ini":   ("code", "text/plain"),
    ".cfg":   ("code", "text/plain"),
    ".conf":  ("code", "text/plain"),
    ".sql":   ("code", "text/x-sql"),
    # Images — described via LLM vision API at ingest time
    ".jpg":  ("image", "image/jpeg"),
    ".jpeg": ("image", "image/jpeg"),
    ".png":  ("image", "image/png"),
    ".bmp":  ("image", "image/bmp"),
    ".heic": ("image", "image/heic"),
}

# Derived views — kept as module-level constants for O(1) membership tests.
SUPPORTED_SUFFIXES: frozenset[str] = frozenset(_FILE_TYPE_INFO)

# Extensionless dotfiles / config files that should always be ingested.
# Matched against the file's name (including a leading dot) case-insensitively.
SUPPORTED_DOTFILES: frozenset[str] = frozenset(
    {
        ".gitconfig", ".gitignore", ".npmrc", ".yarnrc",
        ".bashrc", ".bash_profile", ".zshrc", ".profile",
        ".env", ".envrc",
        "hosts",              # C:\Windows\System32\drivers\etc\hosts
        "makefile", "dockerfile",
    }
)

# Suffixes whose text can be read directly without a format library.
_PLAIN_TEXT_SUFFIXES: frozenset[str] = frozenset(
    ext for ext, (rtype, _) in _FILE_TYPE_INFO.items() if rtype == "code"
) | frozenset({".md", ".txt", ".ini", ".cfg", ".conf"})

# Image suffixes handled via LLM vision description.
_IMAGE_SUFFIXES: frozenset[str] = frozenset({".jpg", ".jpeg", ".png", ".bmp", ".heic"})

_IMAGE_MIME: dict[str, str] = {
    ".jpg":  "image/jpeg",
    ".jpeg": "image/jpeg",
    ".png":  "image/png",
    ".bmp":  "image/bmp",
    ".heic": "image/jpeg",   # HEIC is re-encoded to JPEG bytes before sending
}

_IMAGE_VISION_PROMPT = (
    "Describe this image in detail. Include every visible object, person, "
    "text, scene, activity, color, and any other notable feature. "
    "Your description will be used as searchable text — be specific and thorough."
)


def _is_supported(path: Path) -> bool:
    """Return True if *path* should be ingested by the adapter."""
    name_lower = path.name.lower()
    return (
        path.suffix.lower() in SUPPORTED_SUFFIXES
        or name_lower in SUPPORTED_DOTFILES
    )


# Per-process cache: base_url → True (vision ok) | False (no vision)
_vision_supported_cache: dict[str, bool] = {}


def _is_vision_supported(base_url: str, timeout: int = 5) -> bool:
    """Return True if the running llama-server has a multimodal projector loaded.

    Calls ``GET <base_url>/props`` and checks the ``multimodal`` field.
    Result is cached for the lifetime of the process so scanning a folder
    only pays the HTTP round-trip once.
    """
    import json as _json
    import urllib.request as _req

    if base_url in _vision_supported_cache:
        return _vision_supported_cache[base_url]

    url = base_url.rstrip("/") + "/props"
    try:
        with _req.urlopen(url, timeout=timeout) as resp:  # noqa: S310
            data = _json.loads(resp.read().decode("utf-8"))
        supported = bool(data.get("multimodal", False))
    except Exception:
        # /props endpoint missing (older llama-server) or server unreachable —
        # assume no vision so we don't block ingest.
        supported = False

    _vision_supported_cache[base_url] = supported
    if not supported:
        logger.info(
            "llama-server at %s does not report multimodal support — "
            "image files will be ingested without a vision description. "
            "To enable: add --mmproj <projector.gguf> to the server launch command "
            "and set llama_cpp.mmproj_path in egovault.toml.",
            base_url,
        )
    return supported


def _describe_image(path: Path) -> str:
    """Return a text description of *path* by calling the LLM vision API.

    Sends the image as a base64 data-URI in the OpenAI multimodal message
    format (``image_url`` content part).  HEIC images are converted to JPEG
    first using ``pillow-heif`` (install with ``pip install pillow-heif``).

    Returns an empty string on any failure so ingest can continue.
    """
    import base64 as _b64
    from egovault.config import get_settings
    from egovault.utils.llm import call_llm_chat

    settings = get_settings()
    if not _is_vision_supported(settings.llm.base_url):
        return ""

    suffix = path.suffix.lower()
    try:
        if suffix == ".heic":
            try:
                from pillow_heif import register_heif_opener  # type: ignore[import-untyped]
                register_heif_opener()
            except ImportError:
                logger.warning(
                    "pillow-heif not installed — cannot describe HEIC image %s. "
                    "Install with: pip install pillow-heif",
                    path.name,
                )
                return "[HEIC image — install pillow-heif to enable vision description]"
            import io
            from PIL import Image  # type: ignore[import-untyped]
            with Image.open(path) as img:
                buf = io.BytesIO()
                img.convert("RGB").save(buf, format="JPEG")
                image_bytes = buf.getvalue()
            mime = "image/jpeg"
        else:
            image_bytes = path.read_bytes()
            mime = _IMAGE_MIME.get(suffix, "image/jpeg")

        llm = settings.llm
        b64 = _b64.b64encode(image_bytes).decode("ascii")
        messages = [
            {
                "role": "user",
                "content": [
                    {"type": "text", "text": _IMAGE_VISION_PROMPT},
                    {"type": "image_url", "image_url": {"url": f"data:{mime};base64,{b64}"}},
                ],
            }
        ]
        return call_llm_chat(
            base_url=llm.base_url,
            model=llm.model,
            messages=messages,
            timeout=llm.timeout_seconds,
            provider=llm.provider,
            api_key=llm.api_key,
        )
    except Exception as exc:
        logger.warning("Vision description failed for %s: %s", path.name, exc)
        return ""


def _extract_pdf_liteparse(path: Path) -> str:
    try:
        from liteparse import LiteParse  # type: ignore[import-untyped]
    except ImportError:
        return ""
    try:
        result = LiteParse().parse(str(path), ocr_enabled=True)
        return result.text or ""
    except Exception:
        return ""


def _extract_text(path: Path) -> str:
    """Extract plain text from *path* using the appropriate library."""
    suffix = path.suffix.lower()

    # Extensionless dotfiles — read as plain text
    if not suffix and path.name.lower() in SUPPORTED_DOTFILES:
        return path.read_text(encoding="utf-8", errors="replace")

    if suffix in (".md", ".txt"):
        return path.read_text(encoding="utf-8", errors="replace")

    # Source code and plain-text config/data formats — read directly
    if suffix in _PLAIN_TEXT_SUFFIXES:
        return path.read_text(encoding="utf-8", errors="replace")

    if suffix == ".html":
        try:
            from bs4 import BeautifulSoup  # type: ignore[import-untyped]
        except ImportError:
            return path.read_text(encoding="utf-8", errors="replace")
        soup = BeautifulSoup(
            path.read_text(encoding="utf-8", errors="replace"), "html.parser"
        )
        return soup.get_text(separator="\n", strip=True)

    if suffix == ".pdf":
        try:
            import pdfplumber  # type: ignore[import-untyped]
        except ImportError:
            return _extract_pdf_liteparse(path)
        texts: list[str] = []
        with pdfplumber.open(path) as pdf:
            for page in pdf.pages:
                t = page.extract_text()
                if t:
                    texts.append(t)
        text = "\n".join(texts)
        # Scanned PDFs have no text layer — fall back to LiteParse OCR.
        if not text.strip():
            text = _extract_pdf_liteparse(path)
        return text

    if suffix == ".docx":
        try:
            import docx  # type: ignore[import-untyped]
        except ImportError:
            return ""
        doc = docx.Document(str(path))
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())

    if suffix == ".epub":
        try:
            import ebooklib  # type: ignore[import-untyped]
            from ebooklib import epub
            from bs4 import BeautifulSoup  # type: ignore[import-untyped]
        except ImportError:
            return ""
        book = epub.read_epub(str(path))
        parts: list[str] = []
        for item in book.get_items_of_type(ebooklib.ITEM_DOCUMENT):
            soup = BeautifulSoup(item.get_content(), "html.parser")
            parts.append(soup.get_text(separator="\n", strip=True))
        return "\n".join(parts)

    if suffix == ".xlsx":
        try:
            import pandas as pd  # type: ignore[import-untyped]
        except ImportError:
            return ""
        df_map: dict = pd.read_excel(path, sheet_name=None)
        parts = []
        for sheet_name, sheet_df in df_map.items():
            parts.append(f"Sheet: {sheet_name}\n{sheet_df.to_string(index=False)}")
        return "\n\n".join(parts)

    if suffix == ".pptx":
        try:
            from pptx import Presentation  # type: ignore[import-untyped]
        except ImportError:
            return ""
        prs = Presentation(str(path))
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text)
        return "\n".join(texts)

    if suffix in _IMAGE_SUFFIXES:
        return _describe_image(path)

    return ""


@register
class LocalInboxAdapter(BasePlatformAdapter):
    """Adapter for a user-managed local file drop folder."""

    platform_id = "local"

    def __init__(self, store: "VaultStore | None" = None) -> None:
        self._store = store

    @classmethod
    def can_handle(cls, source_path: Path) -> bool:
        if not source_path.is_dir():
            return False
        return next(
            (
                f for f in source_path.rglob("*")
                if f.is_file() and _is_supported(f)
            ),
            None,
        ) is not None

    def ingest(self, source_path: Path) -> Iterator[NormalizedRecord]:
        for file_path in sorted(source_path.rglob("*")):
            if not file_path.is_file():
                continue
            if not _is_supported(file_path):
                continue

            stat = file_path.stat()
            mtime_f: float = stat.st_mtime
            size: int = stat.st_size
            fid = compute_file_id(str(file_path), mtime_f, size)

            # Stage 1 — fast metadata check (no file read)
            if self._store is not None and self._store.is_file_known(fid):
                # Re-extract if the previously stored body was empty (e.g. pdfplumber
                # returned nothing on first scan but succeeds now).
                if not self._store.record_needs_body_update(str(file_path)):
                    continue
                # Fall through to re-extract the body and yield a fresh record.

            # Stage 2 — content hash check (catches renames / re-touched files)
            try:
                content_hash = compute_content_hash(file_path)
            except (PermissionError, OSError):
                continue
            if self._store is not None and self._store.is_file_known_by_hash(content_hash):
                # Same empty-body check for the hash-matched path
                if not self._store.record_needs_body_update(str(file_path)):
                    continue

            try:
                body = _extract_text(file_path)
            except (PermissionError, OSError):
                continue
            suffix = file_path.suffix.lower()
            # Dotfiles have no suffix — treat as config/code
            effective_suffix = suffix if suffix else ".conf"
            ts = datetime.fromtimestamp(mtime_f, tz=timezone.utc)

            # Relative path from the scan root — used as thread_name so the
            # LLM sees the full folder structure (e.g. "arduino/blink/blink.ino")
            try:
                rel_path = file_path.relative_to(source_path)
            except ValueError:
                rel_path = Path(file_path.name)

            # Prepend the absolute path as a header so it is both visible in
            # vault context and fully indexed by FTS5 (path components become
            # searchable tokens: "arduino", "blink", etc.)
            path_header = f"File: {file_path}\n---\n"
            body = path_header + body

            record = NormalizedRecord(
                platform="local",
                record_type=_FILE_TYPE_INFO.get(effective_suffix, ("document", "text/plain"))[0],
                timestamp=ts,
                sender_id="user",
                sender_name="user",
                thread_id=str(file_path.parent),
                thread_name=str(rel_path),
                body=body,
                attachments=[str(file_path)],
                raw={
                    "file_name": file_path.name,
                    "file_id": fid,
                    "content_hash": content_hash,
                    "size_bytes": size,
                    "mtime": mtime_f,
                },
                file_path=str(file_path),
                mime_type=_FILE_TYPE_INFO.get(effective_suffix, ("document", "text/plain"))[1],
            )
            yield record
